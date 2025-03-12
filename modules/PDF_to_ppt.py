import streamlit as st
import PyPDF2
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_bytes
import io
import tempfile
import os
import zipfile
from utils.common import return_to_main
import base64
import shutil
from utils.common import cleanup_temp_dirs

def pdf_to_pptx(pdf_file, progress_placeholder=None):
    # 创建一个新的PPT演示文稿
    prs = Presentation()
    
    # 读取PDF文件
    pdf_bytes = pdf_file.getvalue()
    pdf_reader = PyPDF2.PdfReader(io.BytesIO(pdf_bytes))
    
    # 获取PDF页数
    num_pages = len(pdf_reader.pages)
    
    # 使用临时目录存储图像
    with tempfile.TemporaryDirectory() as path:
        # 将PDF页面转换为图像
        images = convert_from_bytes(pdf_bytes, dpi=300, output_folder=path)
        
        # 为每个PDF页面创建一个PPT幻灯片
        for i, image in enumerate(images):
            # 更新单个文件转换进度
            if progress_placeholder:
                progress_placeholder.progress((i + 1) / num_pages)
            
            # 添加一个空白幻灯片
            slide_layout = prs.slide_layouts[6]  # 使用空白布局
            slide = prs.slides.add_slide(slide_layout)
            
            # 保存图像到临时文件
            img_path = os.path.join(path, f"slide_{i}.png")
            image.save(img_path, "PNG")
            
            # 将图像添加到幻灯片
            left = top = Inches(0)
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            
            # 添加图片并调整大小以适应幻灯片
            pic = slide.shapes.add_picture(img_path, left, top, width=slide_width, height=slide_height)
    
    # 将演示文稿保存到内存中
    pptx_io = io.BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    
    return pptx_io, num_pages

def create_zip_file(pptx_files):
    """将多个PPT文件打包成zip文件"""
    zip_io = io.BytesIO()
    with zipfile.ZipFile(zip_io, 'w') as zip_file:
        for filename, pptx_data in pptx_files:
            zip_file.writestr(filename, pptx_data.getvalue())
    zip_io.seek(0)
    return zip_io

def pdf_to_ppt():
    return_to_main()
    st.title("PDF➡️PPT")
    st.write("上传一个或多个PDF文件，将其转换为PowerPoint演示文稿")
    
    # 初始化临时目录跟踪
    if 'temp_dirs' not in st.session_state:
        st.session_state.temp_dirs = []
    
    
    uploaded_files = st.file_uploader("选择PDF文件", type="pdf", accept_multiple_files=True)
    
    if uploaded_files:
        if st.button("开始转换"):
            # 清理之前的临时目录
            cleanup_temp_dirs()
            
            # 创建新的临时目录
            temp_dir = tempfile.mkdtemp()
            st.session_state.temp_dirs.append(temp_dir)
            
            with st.spinner("正在转换中，请稍候..."):
                try:
                    # 用于存储所有转换后的PPT文件
                    converted_files = []
                    total_pages = 0
                    
                    # 显示进度条
                    progress_bar = st.progress(0)
                    status_text = st.empty()
                    
                    for i, uploaded_file in enumerate(uploaded_files):
                        # 更新总体进度条
                        progress = (i) / len(uploaded_files)
                        progress_bar.progress(progress)
                        
                        # 显示当前处理的文件名
                        file_name = uploaded_file.name
                        status_text.text(f"正在处理: {file_name}")
                        
                        # 转换PDF
                        pptx_io, num_pages = pdf_to_pptx(uploaded_file)
                        total_pages += num_pages
                        
                        # 生成文件名（去除.pdf后缀并添加.pptx）
                        filename = uploaded_file.name.rsplit('.', 1)[0] + '.pptx'
                        converted_files.append((filename, pptx_io))
                        
                        # 更新总体进度条
                        progress = (i + 1) / len(uploaded_files)
                        progress_bar.progress(progress)
                    
                    # 如果只有一个文件，直接提供下载链接
                    if len(converted_files) == 1:
                        filename, pptx_io = converted_files[0]
                        st.success(f"转换成功！PDF有{total_pages}页。")
                        
                        # 将pptx数据转换为base64
                        pptx_data = pptx_io.getvalue()
                        b64_pptx = base64.b64encode(pptx_data).decode()
                        
                        # 添加CSS样式
                        st.markdown("""
                        <style>
                        .download-button {
                            display: inline-block;
                            padding: 12px 24px;
                            background-color: #4CAF50;
                            color: white !important;
                            text-align: center;
                            text-decoration: none;
                            font-size: 16px;
                            font-weight: bold;
                            margin: 15px 5px;
                            border-radius: 8px;
                            cursor: pointer;
                            transition: all 0.3s;
                            box-shadow: 0 4px 6px rgba(0,0,0,0.1);
                            border: none;
                        }
                        .download-button:hover {
                            background-color: #45a049;
                            transform: translateY(-2px);
                            box-shadow: 0 6px 8px rgba(0,0,0,0.15);
                        }
                        .download-button:active {
                            transform: translateY(0);
                            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
                        }
                        </style>
                        """, unsafe_allow_html=True)
                        
                        # 创建HTML下载链接
                        href = f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{b64_pptx}" download="{filename}" class="download-button">下载PPT文件</a>'
                        st.markdown(href, unsafe_allow_html=True)
                    # 如果有多个文件，创建zip包
                    else:
                        zip_io = create_zip_file(converted_files)
                        st.success(f"转换成功！共转换{len(converted_files)}个文件，总计{total_pages}页。")
                        
                        # 将zip数据转换为base64
                        zip_data = base64.b64encode(zip_io.getvalue()).decode()
                        
                        # 创建HTML下载链接
                        href = f'<a href="data:application/zip;base64,{zip_data}" download="converted_presentations.zip" class="download-button">下载ZIP文件</a>'
                        st.markdown(href, unsafe_allow_html=True)
                
                except Exception as e:
                    st.error(f"转换过程中出错: {str(e)}")
